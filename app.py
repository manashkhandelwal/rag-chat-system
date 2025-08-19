import streamlit as st
import pandas as pd
import base64
import os
import tempfile
from datetime import datetime
from pathlib import Path
import zipfile
import io

# Try to use python-magic-bin on Windows, fallback to mimetypes
try:
    import magic  # from python-magic-bin
    MAGIC_AVAILABLE = True
except ImportError:
    import mimetypes
    MAGIC_AVAILABLE = False
    st.warning("⚠️ 'python-magic' not found, falling back to extension/mimetype detection")

# Document processing imports
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
import openpyxl
from openpyxl import load_workbook


# LangChain imports
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores import Chroma
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from langchain.schema import Document

# Alternative document processing (instead of textract)
from alternative_processors import AlternativeDocumentProcessor, display_processor_status

import asyncio

try:
    asyncio.get_running_loop()
except RuntimeError:
    asyncio.set_event_loop(asyncio.new_event_loop())


class FileFormatDetector:
    """Advanced file format detection using multiple methods"""
    
    def __init__(self):
        self.supported_formats = {
            'pdf', 'docx', 'doc', 'pptx', 'ppt', 
            'xlsx', 'xls', 'csv', 'txt'
        }
    
    def detect(self, file_path):
        """Multi-method file format detection"""
        try:
            if MAGIC_AVAILABLE:
                # Method 1: Magic number detection
                mime_type = magic.from_file(file_path, mime=True)
                format_mapping = {
                    'application/pdf': 'pdf',
                    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
                    'application/msword': 'doc',
                    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
                    'application/vnd.ms-powerpoint': 'ppt',
                    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
                    'application/vnd.ms-excel': 'xls',
                    'text/csv': 'csv',
                    'text/plain': 'txt',
                    'application/zip': 'zip'
                }
                
                detected_format = format_mapping.get(mime_type, 'unknown')
                
                # Handle Office docs that appear as ZIP
                if detected_format == 'zip':
                    detected_format = self._handle_zip_office_docs(file_path)
                    
                if detected_format != 'unknown':
                    return detected_format
            else:
                # Fallback: use mimetypes
                mime_type, _ = mimetypes.guess_type(file_path)
                if mime_type:
                    return mime_type.split("/")[-1]
                
        except Exception as e:
            st.warning(f"Magic detection failed: {e}")
        
        # Fallback to extension
        return self._detect_by_extension(file_path)
    
    def _handle_zip_office_docs(self, file_path):
        """Handle Office documents that appear as ZIP files"""
        try:
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                files = zip_ref.namelist()
                if any('word/' in f for f in files):
                    return 'docx'
                elif any('ppt/' in f for f in files):
                    return 'pptx'
                elif any('xl/' in f for f in files):
                    return 'xlsx'
        except Exception:
            pass
        return self._detect_by_extension(file_path)
    
    def _detect_by_extension(self, file_path):
        """Fallback method using file extensions"""
        extension = Path(file_path).suffix.lower()
        extension_mapping = {
            '.pdf': 'pdf', '.docx': 'docx', '.doc': 'doc',
            '.pptx': 'pptx', '.ppt': 'ppt', '.xlsx': 'xlsx',
            '.xls': 'xls', '.csv': 'csv', '.txt': 'txt'
        }
        return extension_mapping.get(extension, 'unknown')
    
    def is_supported(self, file_path):
        """Check if file format is supported"""
        detected_format = self.detect(file_path)
        return detected_format in self.supported_formats

class DocumentProcessor:
    """Enhanced document processing for multiple formats"""
    
    def __init__(self):
        self.detector = FileFormatDetector()
        self.alternative_processor = AlternativeDocumentProcessor()
    
    def has_complex_layout(self, pdf_path):
        """Detect complex layouts in PDFs"""
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(pdf_path)
            complexity_score = 0
            
            for page_num in range(min(3, len(doc))):
                page = doc[page_num]
                text = page.get_text()
                
                # Table indicators
                if text.count('\t') > 5 or text.count('|') > 3:
                    complexity_score += 2
                
                # Images and drawings
                if len(page.get_images()) > 0:
                    complexity_score += 1
                if len(page.get_drawings()) > 5:
                    complexity_score += 2
                
                # Multi-column detection
                blocks = page.get_text("dict").get("blocks", [])
                if len(blocks) > 10:
                    complexity_score += 1
            
            doc.close()
            return complexity_score >= 3
            
        except Exception:
            return False
    
    def extract_text_from_file(self, file_path, file_format=None):
        """Extract text from various file formats"""
        if file_format is None:
            file_format = self.detector.detect(file_path)
        
        try:
            if file_format == 'pdf':
                return self._extract_pdf_text(file_path)
            elif file_format == 'docx':
                return self._extract_docx_text(file_path)
            elif file_format == 'doc':
                return self._extract_doc_text(file_path)
            elif file_format in ['pptx', 'ppt']:
                return self._extract_ppt_text(file_path)
            elif file_format in ['xlsx', 'xls']:
                return self._extract_excel_text(file_path)
            elif file_format == 'csv':
                return self._extract_csv_text(file_path)
            elif file_format == 'txt':
                return self._extract_txt_text(file_path)
            else:
                raise ValueError(f"Unsupported format: {file_format}")
                
        except Exception as e:
            st.error(f"Error processing {file_format} file: {e}")
            return ""
    
    def _extract_pdf_text(self, file_path):
        """Extract text from PDF files"""
        text = ""
        try:
            # Check for complex layout
            if self.has_complex_layout(file_path):
                st.info("Complex PDF detected, using advanced processing...")
                # Try alternative processors first
                alt_text = self.alternative_processor.extract_pdf_text_alternative(file_path)
                if alt_text and alt_text.strip():
                    return alt_text
            
            # Fallback to PyPDF2
            pdf_reader = PdfReader(file_path)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
                    
        except Exception as e:
            st.error(f"Error reading PDF: {e}")
            
        return text
    
    def _extract_docx_text(self, file_path):
        """Extract text from DOCX files"""
        try:
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            return text
        except Exception as e:
            st.error(f"Error reading DOCX: {e}")
            return ""
    
    def _extract_doc_text(self, file_path):
        """Extract text from DOC files using alternative methods"""
        return self.alternative_processor.extract_doc_text(file_path)
    
    def _extract_ppt_text(self, file_path):
        """Extract text from PowerPoint files"""
        try:
            if file_path.endswith('.pptx'):
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            else:
                # For .ppt files, use alternative processor
                return self.alternative_processor.extract_ppt_text(file_path)
        except Exception as e:
            st.error(f"Error reading PowerPoint: {e}")
            return ""
    
    def _extract_excel_text(self, file_path):
        """Extract text from Excel files"""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            text = ""
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text += f"\n--- Sheet: {sheet_name} ---\n"
                
                # Convert DataFrame to text
                for column in df.columns:
                    text += f"{column}: "
                    values = df[column].dropna().astype(str).tolist()
                    text += ", ".join(values[:100])  # Limit values
                    text += "\n"
                    
            return text
        except Exception as e:
            st.error(f"Error reading Excel: {e}")
            return ""
    
    def _extract_csv_text(self, file_path):
        """Extract text from CSV files"""
        try:
            df = pd.read_csv(file_path)
            text = "CSV Data:\n"
            
            # Add column headers
            text += "Columns: " + ", ".join(df.columns.tolist()) + "\n\n"
            
            # Add sample data and summary
            text += "Sample data (first 10 rows):\n"
            text += df.head(10).to_string(index=False)
            
            text += f"\n\nDataset info: {len(df)} rows, {len(df.columns)} columns"
            
            return text
        except Exception as e:
            st.error(f"Error reading CSV: {e}")
            return ""
    
    def _extract_txt_text(self, file_path):
        """Extract text from TXT files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                st.error(f"Error reading TXT file: {e}")
                return ""

class RAGChatSystem:
    """Enhanced RAG Chat System with Chroma vector store"""
    
    def __init__(self):
        self.document_processor = DocumentProcessor()
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
        
    def process_documents(self, uploaded_files, api_key):
        """Process multiple document types and create vector store"""
        all_documents = []
        processed_files = []
        
        with st.spinner("Processing documents..."):
            progress_bar = st.progress(0)
            
            for i, uploaded_file in enumerate(uploaded_files):
                try:
                    # (your file processing logic...)

                    # Ensure event loop exists before embeddings
                    import asyncio
                    try:
                        asyncio.get_running_loop()
                    except RuntimeError:
                        asyncio.set_event_loop(asyncio.new_event_loop())

                    # Create embeddings and vector store
                    embeddings = GoogleGenerativeAIEmbeddings(
                        model="models/embedding-001",
                        google_api_key=api_key
                    )
                    # Save uploaded file temporarily
                    with tempfile.NamedTemporaryFile(delete=False, suffix=f"_{uploaded_file.name}") as tmp_file:
                        tmp_file.write(uploaded_file.getbuffer())
                        tmp_file_path = tmp_file.name
                    
                    # Detect file format
                    file_format = self.document_processor.detector.detect(tmp_file_path)
                    
                    if not self.document_processor.detector.is_supported(tmp_file_path):
                        st.warning(f"Unsupported file format: {uploaded_file.name}")
                        continue
                    
                    # Extract text
                    text = self.document_processor.extract_text_from_file(tmp_file_path, file_format)
                    
                    if text.strip():
                        # Create chunks
                        chunks = self.text_splitter.split_text(text)
                        
                        # Create documents with metadata
                        for j, chunk in enumerate(chunks):
                            doc = Document(
                                page_content=chunk,
                                metadata={
                                    "source": uploaded_file.name,
                                    "file_type": file_format,
                                    "chunk_id": j,
                                    "upload_time": datetime.now().isoformat()
                                }
                            )
                            all_documents.append(doc)
                        
                        processed_files.append({
                            "name": uploaded_file.name,
                            "format": file_format,
                            "chunks": len(chunks),
                            "text_length": len(text)
                        })
                    
                    # Clean up temp file
                    os.unlink(tmp_file_path)
                    
                    # Update progress
                    progress_bar.progress((i + 1) / len(uploaded_files))
                    
                except Exception as e:
                    st.error(f"Error processing {uploaded_file.name}: {e}")
            
            progress_bar.empty()
        
        if all_documents:
            # Create embeddings and vector store
            embeddings = GoogleGenerativeAIEmbeddings(
                model="models/embedding-001",
                google_api_key=api_key
            )
            
            # Create Chroma vector store
            vectorstore = Chroma.from_documents(
                documents=all_documents,
                embedding=embeddings,
                persist_directory="./chroma_db"
            )
            
            # Persist the database
            vectorstore.persist()
            
            # Display processing summary
            st.success(f"✅ Processed {len(processed_files)} files successfully!")
            
            with st.expander("📄 Processing Summary"):
                summary_df = pd.DataFrame(processed_files)
                st.dataframe(summary_df)
                
                total_chunks = sum(f["chunks"] for f in processed_files)
                total_text = sum(f["text_length"] for f in processed_files)
                st.write(f"**Total chunks created:** {total_chunks}")
                st.write(f"**Total text processed:** {total_text:,} characters")
            
            return vectorstore, processed_files
        else:
            st.error("No documents were successfully processed!")
            return None, []
    
    def get_conversational_chain(self, api_key, research_mode=False):
        """Create conversational chain with different prompts for different modes"""
        
        if research_mode:
            # Deep Research Mode Prompt
            prompt_template = """
            You are a research assistant conducting deep analysis. Based on the provided context, give a comprehensive, 
            well-structured response that includes:
            
            1. **Executive Summary**: Key findings in 2-3 sentences
            2. **Detailed Analysis**: Thorough examination of the topic
            3. **Supporting Evidence**: Relevant quotes and data from the documents
            4. **Cross-References**: Connections between different sources
            5. **Implications**: What this means and potential impacts
            6. **Areas for Further Investigation**: What questions remain
            
            Context:\n {context}\n
            Research Question: \n{question}\n
            
            Provide a detailed research response:
            """
        else:
            # Standard Chat Mode Prompt
            prompt_template = """
            Answer the question based on the provided context. Be conversational but accurate.
            If the answer is not in the provided context, say "I don't have that information in the uploaded documents."
            
            Context:\n {context}\n
            Question: \n{question}\n
            
            Answer:
            """
        
        model = ChatGoogleGenerativeAI(
            model="gemini-1.5-flash",
            temperature=0.3 if not research_mode else 0.1,
            google_api_key=api_key
        )
        
        prompt = PromptTemplate(
            template=prompt_template,
            input_variables=["context", "question"]
        )
        
        chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
        return chain
    
    def query_documents(self, question, api_key, research_mode=False, k=4):
        """Query the vector store and get response"""
        try:
            # Load existing vector store
            embeddings = GoogleGenerativeAIEmbeddings(
                model="models/embedding-001",
                google_api_key=api_key
            )
            
            vectorstore = Chroma(
                persist_directory="./chroma_db",
                embedding_function=embeddings
            )
            
            # Perform similarity search
            if research_mode:
                # Get more documents for research mode
                docs = vectorstore.similarity_search(question, k=k*2)
            else:
                docs = vectorstore.similarity_search(question, k=k)
            
            if not docs:
                return "No relevant information found in the uploaded documents.", []
            
            # Get conversational chain
            chain = self.get_conversational_chain(api_key, research_mode)
            
            # Get response
            response = chain(
                {"input_documents": docs, "question": question},
                return_only_outputs=True
            )
            
            # Extract source information
            sources = []
            for doc in docs:
                sources.append({
                    "source": doc.metadata.get("source", "Unknown"),
                    "file_type": doc.metadata.get("file_type", "Unknown"),
                    "chunk_id": doc.metadata.get("chunk_id", 0),
                    "content_preview": doc.page_content[:200] + "..."
                })
            
            return response['output_text'], sources
            
        except Exception as e:
            st.error(f"Error querying documents: {e}")
            return "Sorry, I encountered an error while processing your question.", []

def main():
    st.set_page_config(
        page_title="Advanced RAG Chat System",
        page_icon="🤖",
        layout="wide"
    )
    
    st.title("🤖 Advanced RAG Chat System")
    st.markdown("*Chat with your documents using AI - Supports PDF, Word, PowerPoint, Excel, CSV, and more!*")
    
    # Initialize session state
    if 'conversation_history' not in st.session_state:
        st.session_state.conversation_history = []
    if 'rag_system' not in st.session_state:
        st.session_state.rag_system = RAGChatSystem()
    if 'documents_processed' not in st.session_state:
        st.session_state.documents_processed = False
    
    # Sidebar configuration
    with st.sidebar:
        st.header("🔧 Configuration")
        
        # API Key input
        api_key = st.text_input("🔑 Google API Key:", type="password")
        if not api_key:
            st.warning("Please enter your Google API Key")
            st.markdown("[Get API Key](https://ai.google.dev/)")
            return
        
        # Display processor status
        display_processor_status()
        
        st.divider()
        
        # File upload
        st.subheader("📂 Upload Documents")
        uploaded_files = st.file_uploader(
            "Choose files",
            accept_multiple_files=True,
            type=['pdf', 'docx', 'doc', 'pptx', 'ppt', 'xlsx', 'xls', 'csv', 'txt']
        )
        
        if uploaded_files:
            st.write(f"📁 {len(uploaded_files)} files selected")
            
            if st.button("🚀 Process Documents", type="primary"):
                vectorstore, processed_files = st.session_state.rag_system.process_documents(
                    uploaded_files, api_key
                )
                if vectorstore:
                    st.session_state.documents_processed = True
                    st.session_state.processed_files = processed_files
        
        st.divider()
        
        # Controls
        st.subheader("🎛️ Controls")
        col1, col2 = st.columns(2)
        
        with col1:
            if st.button("🔄 Reset Chat"):
                st.session_state.conversation_history = []
                st.rerun()
        
        with col2:
            if st.button("🗑️ Clear All"):
                st.session_state.conversation_history = []
                st.session_state.documents_processed = False
                if os.path.exists("./chroma_db"):
                    import shutil
                    shutil.rmtree("./chroma_db")
                st.rerun()
        
        # Export chat history
        if st.session_state.conversation_history:
            st.divider()
            st.subheader("💾 Export")
            
            df = pd.DataFrame(st.session_state.conversation_history)
            csv = df.to_csv(index=False)
            b64 = base64.b64encode(csv.encode()).decode()
            href = f'<a href="data:file/csv;base64,{b64}" download="chat_history.csv">Download Chat History</a>'
            st.markdown(href, unsafe_allow_html=True)
    
    # Main chat interface
    if not st.session_state.documents_processed:
        st.info("👆 Please upload and process documents to start chatting!")
        return
    
    # Chat mode selection
    col1, col2 = st.columns([3, 1])
    with col1:
        user_question = st.text_input("💬 Ask a question about your documents:")
    with col2:
        research_mode = st.checkbox("🔬 Deep Research Mode", help="Get comprehensive analysis")
    
    # Process question
    if user_question:
        with st.spinner("🤔 Thinking..."):
            response, sources = st.session_state.rag_system.query_documents(
                user_question, api_key, research_mode
            )
        
        # Add to conversation history
        st.session_state.conversation_history.append({
            'timestamp': datetime.now(),
            'question': user_question,
            'answer': response,
            'mode': 'Research' if research_mode else 'Chat',
            'sources': len(sources)
        })
        
        # Display response
        st.markdown("### 🤖 Assistant Response")
        st.markdown(response)
        
        # Display sources
        if sources:
            with st.expander(f"📚 Sources ({len(sources)} documents)"):
                for i, source in enumerate(sources):
                    st.write(f"**{i+1}. {source['source']}** ({source['file_type']})")
                    st.write(f"Preview: {source['content_preview']}")
                    st.divider()
    
    # Display chat history
    if st.session_state.conversation_history:
        st.markdown("### 💬 Chat History")
        
        for i, chat in enumerate(reversed(st.session_state.conversation_history)):
            with st.container():
                col1, col2 = st.columns([4, 1])
                with col1:
                    st.markdown(f"**Q:** {chat['question']}")
                with col2:
                    st.caption(f"{chat['mode']} • {chat['timestamp'].strftime('%H:%M')}")
                
                st.markdown(f"**A:** {chat['answer'][:500]}{'...' if len(chat['answer']) > 500 else ''}")
                
                if i < len(st.session_state.conversation_history) - 1:
                    st.divider()

if __name__ == "__main__":
    main()